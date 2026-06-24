# -*- coding: utf-8 -*-
"""Builds the Zomato Group Ordering PM case-study deck (.pptx).
Lighter / minimal theme, Zomato-red accent. Reuses the ZEE deck engine.
Currently: slides 1-2 (Cover + The Launch / History)."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ---------- palette ----------
INK   = RGBColor(0x16,0x1B,0x26)
INK2  = RGBColor(0x5B,0x65,0x77)
INK3  = RGBColor(0x9A,0xA2,0xB0)
LINE  = RGBColor(0xE3,0xE7,0xEF)
WHITE = RGBColor(0xFF,0xFF,0xFF)
SOFT  = RGBColor(0xF6,0xF8,0xFB)
ZRED  = RGBColor(0xE2,0x37,0x44)   # Zomato primary red
ZDEEP = RGBColor(0xCB,0x20,0x2D)   # deeper red
DARK  = RGBColor(0x14,0x1A,0x28)
CRIM  = RGBColor(0xC8,0x36,0x2F)
STEEL = RGBColor(0x25,0x63,0xA8)
PITCH = RGBColor(0x1A,0x7A,0x4A)
VIOL  = RGBColor(0x6D,0x4A,0xA8)
GOLD  = RGBColor(0xB0,0x82,0x2B)
AMBER = RGBColor(0xC9,0x8A,0x12)
ZBG     = RGBColor(0xFD,0xEC,0xEA)
ZBD     = RGBColor(0xF6,0xCD,0xC9)
FAILBG  = RGBColor(0xFC,0xED,0xEC)
FAILBD  = RGBColor(0xF3,0xC8,0xC5)
FIXBG   = RGBColor(0xE9,0xF5,0xEE)
FIXBD   = RGBColor(0xBF,0xE3,0xCD)
STEELBG = RGBColor(0xEA,0xF1,0xFA)
STEELBD = RGBColor(0xC6,0xDA,0xF0)
GOLDBG  = RGBColor(0xF8,0xF1,0xE0)
GOLDBD  = RGBColor(0xE7,0xD5,0xA8)
FONT = "Inter"
FONTD = "Archivo"

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]

def slide():
    return prs.slides.add_slide(BLANK)

def _set_radius(shape, r=0.08):
    try: shape.adjustments[0] = r
    except Exception: pass

def rect(s, x, y, w, h, fill=WHITE, line=None, line_w=1.0, radius=0.08):
    shp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
                             Inches(x), Inches(y), Inches(w), Inches(h))
    if radius: _set_radius(shp, radius)
    if fill is None: shp.fill.background()
    else: shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None: shp.line.fill.background()
    else: shp.line.color.rgb = line; shp.line.width = Pt(line_w)
    shp.shadow.inherit = False
    return shp

def bar(s, x, y, w, h, fill):
    shp = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    shp.line.fill.background(); shp.shadow.inherit = False
    return shp

def text(s, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
         space_after=2, line_spacing=1.0, wrap=True):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space_after); p.space_before = Pt(0)
        p.line_spacing = line_spacing
        for r in para:
            run = p.add_run(); run.text = r.get("t","")
            f = run.font
            f.size = Pt(r.get("sz",12)); f.bold = r.get("b",False)
            f.name = r.get("fn",FONT); f.color.rgb = r.get("c",INK)
            f.italic = r.get("it",False)
            sp = r.get("sp")
            if sp is not None: run.font._rPr.set('spc', str(int(sp)))
    return tb

def P(*runs): return list(runs)
def R(t, sz=12, c=INK, b=False, fn=FONT, sp=None, it=False):
    return {"t":t,"sz":sz,"c":c,"b":b,"fn":fn,"sp":sp,"it":it}

def card(s, x, y, w, h, label, lcolor, heading, body, fill=WHITE, bd=LINE,
         hsize=15, bsize=11.5, accent_top=None):
    rect(s, x, y, w, h, fill=fill, line=bd, line_w=1.0, radius=0.06)
    if accent_top: bar(s, x+0.12, y, w-0.24, 0.055, accent_top)
    pad = 0.16; paras = []
    if label:   paras.append(P(R(label.upper(), 9.5, lcolor, True, FONT, 60)))
    if heading: paras.append(P(R(heading, hsize, INK, True, FONTD)))
    if body:    paras.append(P(*body) if isinstance(body, list) else P(R(body, bsize, INK2)))
    text(s, x+pad, y+pad-0.02, w-2*pad, h-2*pad, paras, space_after=4, line_spacing=1.06)

def footer_links(s, items, page):
    bar(s, 0, 7.16, 13.333, 0.012, LINE)
    tb = s.shapes.add_textbox(Inches(0.7), Inches(7.24), Inches(11.0), Inches(0.26))
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left=0; tf.margin_right=0; tf.margin_top=0; tf.margin_bottom=0
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
    lab = p.add_run(); lab.text = "SOURCES  "
    lab.font.size=Pt(8.5); lab.font.bold=True; lab.font.name=FONT; lab.font.color.rgb=INK3
    lab.font._rPr.set('spc','40')
    for i,(name,url) in enumerate(items):
        if i>0:
            sep=p.add_run(); sep.text="   ·   "; sep.font.size=Pt(9); sep.font.color.rgb=INK3; sep.font.name=FONT
        r=p.add_run(); r.text=name; r.font.size=Pt(9); r.font.name=FONT; r.font.color.rgb=STEEL
        try: r.hyperlink.address=url
        except Exception: pass
    text(s, 12.0, 7.24, 0.9, 0.24, [P(R(page, 9, INK3, True, FONT, 40))],
         align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)

def chip(s, x, y, w, h, label, fill, tc=WHITE):
    rect(s, x, y, w, h, fill=fill, line=None, radius=0.45)
    text(s, x, y, w, h, [P(R(label, 9.5, tc, True, FONT, 30))],
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

def arrow(s, x, y, w, h, color=INK3, glyph="→", sz=20):
    text(s, x, y, w, h, [P(R(glyph, sz, color, True))],
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

def flownode(s, x, y, w, h, top, sub, fill=WHITE, bd=LINE, topc=INK):
    rect(s, x, y, w, h, fill=fill, line=bd, radius=0.08)
    runs=[P(R(top, 11.5, topc, True, FONTD))]
    if sub: runs.append(P(R(sub, 9.5, INK2)))
    text(s, x+0.12, y, w-0.24, h, runs, align=PP_ALIGN.CENTER,
         anchor=MSO_ANCHOR.MIDDLE, space_after=2, line_spacing=0.98)

def header(s, num, eyebrow, accent, title_runs, subtitle=None):
    bar(s, 0, 0, 13.333, 7.5, WHITE)
    bar(s, 0, 0, 13.333, 0.16, accent)
    rect(s, 0.7, 0.5, 0.52, 0.4, fill=accent, line=None, radius=0.18)
    text(s, 0.7, 0.5, 0.52, 0.4, [P(R(num, 13, WHITE, True, FONTD))],
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text(s, 1.35, 0.5, 8.4, 0.4, [P(R(eyebrow.upper(), 12, accent, True, FONT, 140))],
         anchor=MSO_ANCHOR.MIDDLE)
    text(s, 10.3, 0.5, 2.35, 0.4,
         [P(R("ZOMATO ", 12.5, INK2, True, FONTD), R("GROUP ORDERING", 12.5, ZRED, True, FONTD))],
         align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
    text(s, 0.7, 0.98, 12.0, 0.62, [title_runs], line_spacing=0.98)
    if subtitle:
        text(s, 0.7, 1.62, 12.0, 0.42, [subtitle], line_spacing=1.05)

def statbox(s, x, y, w, h, num, label, col):
    rect(s, x, y, w, h, fill=SOFT, line=LINE, radius=0.07)
    text(s, x+0.22, y+0.16, w-0.42, h-0.28,
         [P(R(num, 19, col, True, FONTD)), P(R(label, 9.5, INK2))],
         space_after=3, line_spacing=1.04)

def tnode(s, x, y, w, h, date, datecol, title, desc, fill, bd, tc):
    rect(s, x, y, w, h, fill=fill, line=bd, radius=0.06)
    chip(s, x+0.14, y+0.12, 1.55, 0.3, date, datecol)
    text(s, x+0.16, y+0.5, w-0.3, h-0.62,
         [P(R(title, 11.5, tc, True, FONTD)), P(R(desc, 9.3, INK2))],
         space_after=2, line_spacing=1.0)

def pageno(s, t):
    bar(s, 0, 7.16, 13.333, 0.012, LINE)
    text(s, 12.0, 7.24, 0.9, 0.24, [P(R(t, 9, INK3, True, FONT, 40))],
         align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)

def safe_save(prs, name):
    import os
    try:
        prs.save(name); return name
    except PermissionError:
        base, ext = os.path.splitext(name); i = 2
        while True:
            alt = "%s-v%d%s" % (base, i, ext)
            try:
                prs.save(alt)
                print("NOTE: '%s' was locked — saved to '%s' instead." % (name, alt))
                return alt
            except PermissionError:
                i += 1
                if i > 9: raise

# ---------- source registry ----------
SRC = {
 "swiggy_li":("Inc42","https://inc42.com/buzz/zomato-rolls-out-new-feature-for-users-to-place-orders-in-groups/"),
 "swiggy_blog":("Swiggy Engineering","https://blog.swiggy.com/life-at-swiggy/heres-how-swiggy-built-the-group-ordering-feature/"),
 "zomato_x":("Free Press Journal","https://www.freepressjournal.in/business/no-more-passing-the-phone-zomato-boss-deepinder-goyal-announces-new-group-ordering-feature"),
 "zomato_isn":("Indian Startup News","https://indianstartupnews.com/news/deepinder-goyal-led-zomato-launches-group-ordering-feature-6866846"),
 "eternal_fy25":("MediaNama","https://www.medianama.com/2025/07/223-eternal-fy25-annual-report-highlights-zomato-blinkit/"),
 "eternal_pdf":("Eternal FY25 Results","https://b.zmtcdn.com/investor-relations/d9c290cd23764a09789769c39682276a_1746094084.pdf"),
 "startupstory":("Startup Story","https://startupstorymedia.com/insights-after-swiggy-zomato-launches-group-ordering-feature-all-you-need-to-know/"),
}
def srcs(*keys): return [SRC[k] for k in keys]

# ============================================================ SLIDE 1 · COVER
s = slide()
bar(s, 0, 0, 13.333, 7.5, WHITE)
bar(s, 0, 0, 13.333, 0.16, ZRED)
bar(s, 0.7, 1.02, 0.62, 0.07, ZRED)

text(s, 0.7, 0.62, 11, 0.34,
     [P(R("ZOMATO  ·  PRODUCT TEARDOWN & PRODUCT DESIGN  ·  2026", 12.5, INK2, True, FONT, 110))])

text(s, 0.66, 1.18, 12.2, 1.9,
     [P(R("Group ", 50, INK, True, FONTD), R("Ordering", 50, ZRED, True, FONTD)),
      P(R("re-architected", 50, INK, True, FONTD))],
     line_spacing=0.98)

text(s, 0.7, 3.28, 11.7, 1.0,
     [P(R("Both Swiggy and Zomato shipped the cart. Neither shipped the ", 18, INK2),
        R("coordination", 18, INK, True),
        R(" — so the group went back to WhatsApp. Here's the feature that makes it stick.", 18, INK2))],
     line_spacing=1.18)

stats = [("3.7×","Group-order AOV vs solo · Swiggy launch weekend", ZRED),
         ("6%","of all Swiggy orders, first weekend live", ZDEEP),
         ("₹453","Zomato solo AOV · FY25 baseline", STEEL),
         ("0","platforms that solved split-pay fairness", VIOL)]
cw, gap, x0, top, ch = 2.83, 0.21, 0.7, 4.5, 1.18
for i,(n,l,col) in enumerate(stats):
    x = x0 + i*(cw+gap)
    rect(s, x, top, cw, ch, fill=SOFT, line=LINE, radius=0.07)
    text(s, x+0.24, top+0.18, cw-0.42, ch-0.3,
         [P(R(n, 24, col, True, FONTD)), P(R(l, 10, INK2))], space_after=4, line_spacing=1.05)

bar(s, 0.7, 6.06, 11.93, 0.012, LINE)
text(s, 0.7, 6.18, 9, 0.34,
     [P(R("by ", 13.5, INK2), R("Prakhar Beniwal", 13.5, INK, True),
        R("  ·  Aspiring Product Manager", 13.5, INK2))])
text(s, 0.7, 6.6, 12, 0.5,
     [P(R("Portfolio case study — competitor figures are self-reported launch-window or single-sourced where noted; financials are directional models tagged throughout.", 10, INK3))],
     line_spacing=1.15)
footer_links(s, srcs("swiggy_li","eternal_fy25","zomato_isn"), "01 / 14")

# ===================================== SLIDE 2 · THE LAUNCH / HISTORY
s = slide()
header(s, "02", "The Launch · August 2024", ZRED,
       P(R("Two giants, ", 27, INK, True, FONTD), R("the same idea", 27, ZRED, True, FONTD),
         R(", the same month — ", 27, INK, True, FONTD), R("and the same blind spot", 27, ZDEEP, True, FONTD)),
       P(R("Group ordering wasn't a moonshot — it was a fast-follow race, shipped in weeks off a thin idea. To see what was missing, first read exactly what shipped, and why both rushed.", 12.5, INK2)))

# --- timeline of 4 nodes (red baseline) ---
bar(s, 0.7, 2.04, 11.93, 0.022, ZRED)
tnw, tng, tty, tth = 2.78, 0.27, 1.92, 1.42
txs = [0.7+i*(tnw+tng) for i in range(4)]
tnode(s, txs[0], tty, tnw, tth, "DEC 2023", INK2, "A hackathon castoff",
      "A 2-day prototype — by an engineer who didn't even win — is picked up by Swiggy's PM team. The bet was cheap, not visionary.", SOFT, LINE, INK)
tnode(s, txs[1], tty, tnw, tth, "MID-AUG 2024", STEEL, "Swiggy ships first",
      "QR-first, built for the “party pain point.” One shared cart, host pays; engineered so non-users can join with no app.", STEELBG, STEELBD, STEEL)
tnode(s, txs[2], tty, tnw, tth, "AUG 17, 2024", ZRED, "Zomato clones it in days",
      "Goyal on X: “No more passing the phone around.” A near-identical link-share build, rolled out gradually.", ZBG, ZBD, ZDEEP)
tnode(s, txs[3], tty, tnw, tth, "THAT DAY →", CRIM, "The promise it broke",
      "Asked about splitting the bill, Goyal says “coming soon.” ~22 months later, still no native split-pay.", FAILBG, FAILBD, CRIM)

# --- left: what shipped & where it cracks ---
rect(s, 0.7, 3.5, 7.15, 2.04, fill=WHITE, line=LINE, radius=0.04)
text(s, 0.92, 3.62, 6.8, 0.26, [P(R("WHAT BOTH SHIPPED — AND WHERE IT CRACKS", 10, ZRED, True, FONT, 30))])
bar(s, 0.92, 3.88, 6.71, 0.02, ZBD)
mfx, mfy, mfw, mfh = 0.92, 4.0, 1.5, 0.8
labels = [("Pick","restaurant"),("Share","link / QR"),("Guests add","to one cart"),("Host pays","for everyone")]
fills  = [SOFT, STEELBG, ZBG, FAILBG]
bds    = [LINE, STEELBD, ZBD, FAILBD]
tcs    = [INK, STEEL, ZDEEP, CRIM]
for i,(t1,t2) in enumerate(labels):
    fx = mfx + i*(mfw+0.236)
    flownode(s, fx, mfy, mfw, mfh, t1, t2, fill=fills[i], bd=bds[i], topc=tcs[i])
    if i < 3:
        arrow(s, fx+mfw-0.02, mfy, 0.27, mfh, ZRED, glyph="→", sz=16)
text(s, 0.92, 4.86, 6.8, 0.3,
     [P(R("Mechanically identical: ", 10, INK, True),
        R("host picks a restaurant, shares a link or QR, guests pile into one cart, host checks out for the table.", 10, INK2))],
     line_spacing=1.06)
text(s, 0.92, 5.16, 6.8, 0.34,
     [P(R("Three cracks, baked in day one:  ", 10, ZDEEP, True),
        R("one address · one payer (host silently eats every fee + tip) · zero presence — you can't see who added what.", 10, INK2))],
     line_spacing=1.06)

# --- right: the strategic read (inference) ---
rect(s, 8.0, 3.5, 4.63, 2.04, fill=ZBG, line=ZBD, radius=0.04)
text(s, 8.22, 3.62, 4.2, 0.26, [P(R("THE STRATEGIC READ · WHY THE RUSH", 10, ZDEEP, True, FONT, 20))])
bar(s, 8.22, 3.88, 4.19, 0.02, ZRED)
reads = [
 ("The AOV lever, not love of groups","more mouths = a bigger basket; the 3.7× proved the math overnight. A margin play."),
 ("Cheap to ship","checkout was already solved — group ordering is a thin coordination wrapper. Fast follow."),
 ("A parity race, not a moat","both cloned one mechanic in weeks; whoever cracks coordination + fairness actually wins."),
 ("Built for the launch tweet","single address, host-pays-all, split “coming soon” — the hard retention work was deferred."),
]
ry = 4.0
for lead, detail in reads:
    text(s, 8.22, ry, 4.2, 0.4,
         [P(R("›  "+lead+" — ", 9, ZDEEP, True, FONTD), R(detail, 8.7, INK2))],
         line_spacing=0.98)
    ry += 0.385

# --- proof-of-demand strip ---
rect(s, 0.7, 5.62, 11.93, 0.72, fill=SOFT, line=LINE, radius=0.05)
text(s, 0.92, 5.62, 1.7, 0.72,
     [P(R("PROOF THE", 8.5, ZRED, True, FONT, 20)), P(R("DEMAND IS REAL", 8.5, ZRED, True, FONT, 20))],
     anchor=MSO_ANCHOR.MIDDLE, space_after=0, line_spacing=1.02)
ps = [("3.7×","group AOV vs solo"),("6%","of launch-wknd orders"),
      ("17","real max in one cart"),("72 items","one ₹14,000 order")]
for i,(n,l) in enumerate(ps):
    px = 2.66 + i*2.46
    rect(s, px, 5.72, 2.32, 0.52, fill=WHITE, line=LINE, radius=0.06)
    text(s, px+0.14, 5.74, 2.08, 0.48, [P(R(n, 13.5, ZDEEP, True, FONTD)), P(R(l, 7.8, INK2))],
         space_after=0, line_spacing=0.9)

# --- deep-red cliffhanger ---
rect(s, 0.7, 6.42, 11.93, 0.6, fill=ZDEEP, line=None, radius=0.05)
text(s, 0.95, 6.42, 11.45, 0.6,
     [P(R("→  ", 14, RGBColor(0xFF,0xC9,0xC4), True),
        R("Demand was never the question. Coordination was — ", 12, WHITE, True),
        R("the spike proved the want; the silence after proved the gap.", 12, RGBColor(0xFF,0xE0,0x9E), True))],
     anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0)

footer_links(s, srcs("swiggy_blog","zomato_x","swiggy_li","eternal_fy25"), "02 / 14")

SRC.update({
 "pendo":("Pendo Feature Adoption Report","https://www.pendo.io/resources/the-2019-feature-adoption-report/"),
 "pendo_wral":("WRAL TechWire","https://wraltechwire.com/2020/01/28/pendo-study-with-80-of-features-not-used-software-execs-re-evaluating-success-metrics/"),
})

# ===================================== SLIDE 3 · THE SPIKE THAT FADED
s = slide()
header(s, "03", "The Fade · Where It Leaks", ZRED,
       P(R("A launch spike — ", 27, INK, True, FONTD), R("leaking where it matters most", 27, ZRED, True, FONTD)),
       P(R("The demand is real and proven. The feature still bleeds — and not where you'd think. Walk the funnel: it holds at activation, and breaks at discovery and coordination.", 12.5, INK2)))

# --- leak funnel: 4 stages, arrows, status ---
stages = [
 ("AWARENESS","could group-order","BROAD",   STEEL, STEELBG, STEELBD, None),
 ("DISCOVERY","find the feature","LEAK",      ZRED,  ZBG,     ZBD,     ("↓ LEAK", ZRED)),
 ("ACTIVATION","works when surfaced","HEALTHY",PITCH, FIXBG,   FIXBD,   ("✓ this works", PITCH)),
 ("RETENTION","come back & repeat","THE LEAK", ZDEEP, FAILBG,  FAILBD,  ("↓↓ THE BIG LEAK", ZDEEP)),
]
pw, pgap, py, ph = 2.6, 0.51, 2.24, 1.12
for i,(nm,role,stat,col,fb,bd,leak) in enumerate(stages):
    fx = 0.7 + i*(pw+pgap)
    rect(s, fx, py, pw, ph, fill=fb, line=bd, radius=0.06)
    chip(s, fx+(pw-1.5)/2, py-0.17, 1.5, 0.32, stat, col)
    text(s, fx+0.14, py+0.2, pw-0.28, ph-0.34,
         [P(R(nm, 14, col, True, FONTD)), P(R(role, 9.5, INK2))],
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, space_after=2, line_spacing=1.0)
    if i < 3:
        arrow(s, fx+pw, py, pgap, ph, ZRED, glyph="→", sz=22)
    if leak:
        text(s, fx, py+ph+0.04, pw, 0.3, [P(R(leak[0], 10.5, leak[1], True, FONTD))],
             align=PP_ALIGN.CENTER)

text(s, 0.7, 3.74, 11.93, 0.3,
     [P(R("Activation is fine — the want is proven. ", 11, PITCH, True),
        R("The money bleeds at the two ", 11, INK2), R("ends", 11, ZDEEP, True),
        R(": nobody finds it (discovery), and nobody returns (coordination).", 11, INK2))],
     align=PP_ALIGN.CENTER, line_spacing=1.04)

# --- evidence row: why the spike never became a habit ---
text(s, 0.7, 4.2, 11.93, 0.26, [P(R("WHY THE SPIKE NEVER BECAME A HABIT", 10, ZRED, True, FONT, 40))])
bar(s, 0.7, 4.46, 11.93, 0.02, ZBD)
ev = [(ZBG,ZBD,ZDEEP,"0 in 22 mo","native split-pay shipped — promised at launch, never delivered. The hard part was skipped."),
      (GOLDBG,GOLDBD,AMBER,"~80%","of software features go rarely or never used (Pendo). A buried feature is a dead feature."),
      (STEELBG,STEELBD,STEEL,"0","steady-state group numbers ever reported — only self-reported launch-weekend PR. The silence tells.")]
ecw, eg = 3.84, 0.205
for i,(fb,bd,col,n,l) in enumerate(ev):
    ex = 0.7 + i*(ecw+eg)
    rect(s, ex, 4.58, ecw, 1.28, fill=fb, line=bd, radius=0.06)
    text(s, ex+0.2, 4.72, ecw-0.4, 1.04,
         [P(R(n, 21, col, True, FONTD)), P(R(l, 9.5, INK2))], space_after=4, line_spacing=1.05)
text(s, 0.7, 5.96, 11.93, 0.22,
     [P(R("Spike curve is directional — no platform discloses steady-state group volume; the case rests on the signals above, not a published decline.", 8, INK3, False, FONT, 0, True))])

# --- cliffhanger ---
rect(s, 0.7, 6.4, 11.93, 0.6, fill=ZDEEP, line=None, radius=0.05)
text(s, 0.95, 6.4, 11.45, 0.6,
     [P(R("→  ", 14, RGBColor(0xFF,0xC9,0xC4), True),
        R("A launch is not a habit. To win, Zomato has to remove the reason they leave — ", 12, WHITE, True),
        R("the coordination cost.", 12, RGBColor(0xFF,0xE0,0x9E), True))],
     anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0)
footer_links(s, srcs("zomato_x","pendo","swiggy_li"), "03 / 14")

# ===================================== SLIDE 4 · COORDINATION, NOT CHECKOUT
s = slide()
header(s, "04", "Diagnosis · Coordination, Not Checkout", ZDEEP,
       P(R("Zomato solved checkout. ", 26, INK, True, FONTD),
         R("The real bill is coordination", 26, ZRED, True, FONTD)),
       P(R("The cart was the easy 10%. The unpaid tax is everything around it — five clusters of coordination cost, each feeding the next, most of it borne before and after the app.", 12, INK2)))

# --- reframe: solved vs unsolved ---
rect(s, 0.7, 2.05, 3.95, 0.92, fill=FIXBG, line=FIXBD, radius=0.06)
text(s, 0.9, 2.05, 3.6, 0.92,
     [P(R("SOLVED  ✓", 12, PITCH, True, FONTD)),
      P(R("Checkout & the shared cart — a thin wrapper on an already-solved flow.", 9.2, INK2))],
     anchor=MSO_ANCHOR.MIDDLE, space_after=2, line_spacing=1.0)
arrow(s, 4.68, 2.05, 0.56, 0.92, ZRED, glyph="→", sz=26)
rect(s, 5.32, 2.05, 7.31, 0.92, fill=ZBG, line=ZBD, radius=0.06)
text(s, 5.54, 2.05, 6.95, 0.92,
     [P(R("UNSOLVED  ✕", 12, ZDEEP, True, FONTD),
        R("   Coordination", 12, ZRED, True, FONTD)),
      P(R("The actual job: getting a group from “let's order something” to “it's here, and everyone paid their share.”", 9.2, INK2))],
     anchor=MSO_ANCHOR.MIDDLE, space_after=2, line_spacing=1.0)

# --- the coordination-cost chain (5-cluster table) ---
text(s, 0.7, 3.16, 8.5, 0.26, [P(R("THE COORDINATION-COST CHAIN — EACH FAILURE FEEDS THE NEXT", 10, ZRED, True, FONT, 20))])
chip(s, 11.06, 3.12, 1.57, 0.34, "18 FAILURE POINTS", ZDEEP)
clusters = [
 ("DISCOVERY","find it at all", VIOL, RGBColor(0xF1,0xEC,0xF8),
   ["Buried in the cart","No home-screen entry","No pre-cart consensus"], "MED", 0.5),
 ("COLLECTION","who's in, what they want", STEEL, STEELBG,
   ["Who's even joining?","Phone-passing","What does each want?","Late joiners"], "HIGH", 0.82),
 ("SETTLEMENT","who pays what", ZRED, ZBG,
   ["Host eats all fees","Chasing ₹80s on UPI","“I'll pay later” decay","No split-pay"], "HIGHEST", 1.0),
 ("TIMING","when do we fire", AMBER, GOLDBG,
   ["Waiting on stragglers","Nagging the group","No deadline or nudge"], "HIGH", 0.82),
 ("FULFILMENT","whose food is whose", PITCH, FIXBG,
   ["“Whose biryani?”","Mixed-up bags","One address only","Big-order packaging"], "MED", 0.55),
]
colw, cgap, cy, ch = 2.18, 0.2475, 3.56, 2.18
for i,(nm,sub,col,fb,fails,lvl,frac) in enumerate(clusters):
    cx = 0.7 + i*(colw+cgap)
    rect(s, cx, cy, colw, ch, fill=fb, line=col, line_w=1.0, radius=0.05)
    rect(s, cx, cy, colw, 0.6, fill=col, line=None, radius=0.05)
    text(s, cx+0.12, cy+0.04, colw-0.24, 0.54,
         [P(R(nm, 11, WHITE, True, FONTD)), P(R(sub, 7.5, WHITE)),],
         anchor=MSO_ANCHOR.MIDDLE, space_after=0, line_spacing=0.92)
    text(s, cx+0.13, cy+0.68, colw-0.24, ch-1.08,
         [P(R("• ", 8.5, col, True), R(f, 8.4, INK)) for f in fails],
         space_after=3, line_spacing=0.96)
    # cost meter
    mx, my, mw = cx+0.13, cy+ch-0.32, colw-0.26
    mcol = ZRED if lvl in ("HIGH","HIGHEST") else col
    text(s, mx, my-0.18, mw, 0.16, [P(R("COORD. COST · "+lvl, 6.8, mcol, True, FONT, 10))])
    rect(s, mx, my, mw, 0.11, fill=RGBColor(0xEC,0xEF,0xF4), line=None, radius=0)
    rect(s, mx, my, mw*frac, 0.11, fill=mcol, line=None, radius=0)
    if i < 4:
        arrow(s, cx+colw-0.03, cy+0.13, cgap+0.06, 0.34, ZRED, glyph="→", sz=15)

text(s, 0.7, 5.84, 11.93, 0.3,
     [P(R("The heaviest costs — ", 11, INK2), R("settlement", 11, ZRED, True), R(" and ", 11, INK2),
        R("collection", 11, STEEL, True),
        R(" — are exactly what Zomato's v1 leaves untouched. The cart fixed the cheapest link in the chain.", 11, INK2))],
     align=PP_ALIGN.CENTER, line_spacing=1.04)

# --- cliffhanger -> personas ---
rect(s, 0.7, 6.4, 11.93, 0.6, fill=ZDEEP, line=None, radius=0.05)
text(s, 0.95, 6.4, 11.45, 0.6,
     [P(R("→  ", 14, RGBColor(0xFF,0xC9,0xC4), True),
        R("The coordination tax isn't evenly paid. ", 12, WHITE, True),
        R("Meet the people who pay the most.", 12, RGBColor(0xFF,0xE0,0x9E), True))],
     anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0)
footer_links(s, srcs("swiggy_blog","swiggy_li","zomato_x"), "04 / 14")

def avatar(s, x, y, d, initials, col):
    o = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(d), Inches(d))
    o.fill.solid(); o.fill.fore_color.rgb = col
    o.line.color.rgb = WHITE; o.line.width = Pt(1.75); o.shadow.inherit = False
    text(s, x, y, d, d, [P(R(initials, 12 if len(initials)<3 else 9.5, WHITE, True, FONTD))],
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

def meter(s, x, y, w, frac, col, track=RGBColor(0xEC,0xEF,0xF4)):
    rect(s, x, y, w, 0.11, fill=track, line=None, radius=0)
    rect(s, x, y, max(0.04, w*frac), 0.11, fill=col, line=None, radius=0)

# clusters palette (matches slide 4)
CL = {"DISCOVERY":VIOL, "COLLECTION":STEEL, "SETTLEMENT":ZRED, "TIMING":AMBER, "FULFILMENT":PITCH}

# ===================================== SLIDE 5 · PERSONA · A NIGHT IN THE LIFE
s = slide()
header(s, "05", "Personas · The Coordination Tax", ZRED,
       P(R("11:40 PM. Five hungry flatmates. ", 24, INK, True, FONTD),
         R("One very tired host.", 24, ZRED, True, FONTD)),
       P(R("Meet Rahul and the PG crew — the anchor user. Follow ONE late-night order to feel how heavy the coordination tax really is, even with “group ordering” switched on.", 12, INK2)))

# anchor persona card
rect(s, 0.7, 2.05, 2.7, 3.0, fill=ZBG, line=ZBD, radius=0.05)
avatar(s, 1.7, 2.24, 0.7, "R", ZRED)
chip(s, 1.2, 3.04, 1.7, 0.3, "ANCHOR PERSONA", ZDEEP)
text(s, 0.84, 3.44, 2.42, 0.34, [P(R("Rahul & the PG crew", 13, INK, True, FONTD))], align=PP_ALIGN.CENTER)
text(s, 0.84, 3.82, 2.42, 0.9,
     [P(R("3–8 flatmates · students", 9, INK2)),
      P(R("price-sensitive · late-night", 9, INK2)),
      P(R("UPI + WhatsApp native", 9, INK2)),
      P(R("shared rooms — QR shines", 9, INK2))],
     align=PP_ALIGN.CENTER, space_after=1, line_spacing=1.0)
text(s, 0.84, 4.62, 2.42, 0.42,
     [P(R("WHY THE ANCHOR", 8, ZRED, True, FONT, 30)),
      P(R("Highest pain, densest social graph, biggest viral upside.", 8.7, INK2))],
     align=PP_ALIGN.CENTER, space_after=1, line_spacing=0.98)

# pain journey storyboard (6 steps)
steps5 = [
 ("1","The 40-message debate","20 min of “biryani or pizza?” in the chat. Half never reply.","COLLECTION",0.3),
 ("2","Pass the phone","His phone circles the room; someone deletes an item by mistake.","COLLECTION",0.45),
 ("3","Who's even in?","Two are “still deciding.” The cart sits open for half an hour.","TIMING",0.6),
 ("4","Rahul pays for all","₹2,140 on his card — free delivery only on his Gold.","SETTLEMENT",0.8),
 ("5","The ₹80 chase","Two pay, two “tomorrow,” one ghosts. Day 4: still down ₹430.","SETTLEMENT",1.0),
 ("6","Whose biryani?","One bag, no labels, five identical parcels at the door.","FULFILMENT",0.7),
]
sgx, sgy, scw, sch, sgx2, sgy2 = 3.6, 2.05, 2.85, 1.4, 0.18, 0.2
for i,(num,title,fric,clu,frac) in enumerate(steps5):
    cxx = sgx + (i%3)*(scw+sgx2)
    cyy = sgy + (i//3)*(sch+sgy2)
    ccol = CL[clu]
    rect(s, cxx, cyy, scw, sch, fill=WHITE, line=LINE, radius=0.05)
    avatar(s, cxx+0.12, cyy+0.12, 0.32, num, ccol)
    text(s, cxx+0.52, cyy+0.13, scw-0.62, 0.3, [P(R(title, 10.5, INK, True, FONTD))], anchor=MSO_ANCHOR.MIDDLE)
    text(s, cxx+0.14, cyy+0.48, scw-0.28, 0.5, [P(R(fric, 8.4, INK2))], line_spacing=0.98)
    chip(s, cxx+0.14, cyy+sch-0.34, 1.18, 0.26, clu, ccol)
    meter(s, cxx+1.46, cyy+sch-0.29, scw-1.6, frac, ccol)
    if i%3 < 2 and i != 5:
        arrow(s, cxx+scw-0.02, cyy, sgx2, sch, ZRED, glyph="→", sz=15)

# the real cost
rect(s, 0.7, 5.22, 11.93, 1.0, fill=SOFT, line=LINE, radius=0.05)
text(s, 0.92, 5.3, 11.5, 0.24, [P(R("THE REAL COST OF ONE ORDER", 10, ZRED, True, FONT, 30))])
cost = [("EMOTIONAL", ZDEEP, "The host becomes the group's unpaid project manager — and burns out."),
        ("FINANCIAL", ZRED, "The host is the group's bank: ₹430 still unpaid on day 4."),
        ("BEHAVIOURAL", CRIM, "So he quits the feature: “just pay me, I'll order” — back to square one.")]
for i,(lab,col,txt) in enumerate(cost):
    cx = 0.92 + i*3.92
    rect(s, cx, 5.58, 3.78, 0.56, fill=WHITE, line=LINE, radius=0.06)
    text(s, cx+0.16, 5.62, 3.5, 0.48,
         [P(R(lab+"  ", 8.5, col, True, FONT, 20), R(txt, 8.6, INK2))], anchor=MSO_ANCHOR.MIDDLE, line_spacing=0.96)

# cliffhanger
rect(s, 0.7, 6.42, 11.93, 0.58, fill=ZDEEP, line=None, radius=0.05)
text(s, 0.95, 6.42, 11.45, 0.58,
     [P(R("→  ", 14, RGBColor(0xFF,0xC9,0xC4), True),
        R("And Rahul has it easy — one room, one tap away. ", 12, WHITE, True),
        R("Now scale this to offices, families and 30-person parties.", 12, RGBColor(0xFF,0xE0,0x9E), True))],
     anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0)
footer_links(s, srcs("swiggy_li","zomato_x","eternal_fy25"), "05 / 14")

# ===================================== SLIDE 6 · PERSONA SPECTRUM
s = slide()
header(s, "06", "Personas · It's Not Just Hostels", ZRED,
       P(R("Same tax, ", 26, INK, True, FONTD), R("every table", 26, ZRED, True, FONTD)),
       P(R("Six occasions, one coordination tax — and each group snaps at a different link in the chain. The pain is universal; only the breaking point moves.", 12, INK2)))

people = [
 ("RP","Rahul & PG crew","ANCHOR","3–8 students · late-night · UPI-native",
   "Host is the group's unpaid bank — the ₹80s never come back.","SETTLEMENT",1.0),
 ("OS","Office / WFH squad","","co-located + remote · fixed lunch window",
   "1 PM and half the team still hasn't picked — it fires late, or not at all.","TIMING",0.8),
 ("PH","Party host","","20–30 guests · one giant cart",
   "One person fielding 30 orders in a WhatsApp avalanche (the 72-item cart).","COLLECTION",0.88),
 ("Fam","Family · multi-gen","","mixed app-literacy · one address",
   "Mom dictates her order, someone re-types it; only the host has Gold.","COLLECTION",0.62),
 ("T3","Tier-2/3 friends","","intensely price-sensitive",
   "Every rupee counts — splitting fees by hand kills the whole vibe.","SETTLEMENT",0.9),
 ("GG","Gaming / study crew","","mid-session · same room",
   "Nobody wants to pause to coordinate, so the order just never happens.","TIMING",0.72),
]
pcw, pch, pgx, pgy = 3.84, 1.75, 0.205, 0.22
for i,(ini,nm,tag,crew,pain,clu,frac) in enumerate(people):
    cx = 0.7 + (i%3)*(pcw+pgx)
    cy = 2.12 + (i//3)*(pch+pgy)
    ccol = CL[clu]
    rect(s, cx, cy, pcw, pch, fill=WHITE, line=LINE, radius=0.05)
    avatar(s, cx+0.16, cy+0.16, 0.52, ini, ccol)
    text(s, cx+0.8, cy+0.16, pcw-1.5, 0.52,
         [P(R(nm, 11.5, INK, True, FONTD)), P(R(crew, 8, INK2))], anchor=MSO_ANCHOR.MIDDLE, space_after=0, line_spacing=0.94)
    if tag:
        chip(s, cx+pcw-0.92, cy+0.16, 0.78, 0.28, tag, ZDEEP)
    text(s, cx+0.18, cy+0.8, pcw-0.36, 0.52, [P(R(pain, 9.3, ZDEEP, True))], line_spacing=1.0)
    chip(s, cx+0.18, cy+pch-0.36, 1.2, 0.26, clu, ccol)
    text(s, cx+1.46, cy+pch-0.4, 1.2, 0.2, [P(R("COORD. TAX", 6.8, INK3, True, FONT, 10))])
    meter(s, cx+1.46, cy+pch-0.22, pcw-1.62, frac, ccol)

# summary strip
rect(s, 0.7, 6.0, 11.93, 0.32, fill=ZBG, line=ZBD, radius=0.05)
text(s, 0.7, 6.0, 11.93, 0.32,
     [P(R("6 occasions   ·   18 failure points   ·   ", 11, ZDEEP, True, FONTD),
        R("0 platforms that remove the tax — every group still does the hard part by hand.", 11, INK, True))],
     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# cliffhanger
rect(s, 0.7, 6.42, 11.93, 0.58, fill=ZDEEP, line=None, radius=0.05)
text(s, 0.95, 6.42, 11.45, 0.58,
     [P(R("→  ", 14, RGBColor(0xFF,0xC9,0xC4), True),
        R("The pain is real, universal, and unsolved. ", 12, WHITE, True),
        R("So how has the rest of the world cracked it?", 12, RGBColor(0xFF,0xE0,0x9E), True))],
     anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0)
footer_links(s, srcs("swiggy_li","zomato_x","eternal_fy25"), "06 / 14")

SRC.update({
 "grab":("Inside Grab","https://www.grab.com/inside-grab/stories/grabfood-delivery-group-order/"),
 "spotify_jam":("TechCrunch","https://techcrunch.com/2023/09/26/spotify-launches-jam-a-real-time-collaborative-playlist-controlled-by-up-to-32-people/"),
 "doordash":("DoorDash Help","https://help.doordash.com/work/s/article/Guide-to-Group-Orders"),
 "ubereats":("Uber Eats Help","https://help.uber.com/en/ubereats/restaurants/article/how-to-place-a-group-order-"),
 "meituan":("Rest of World","https://restofworld.org/2024/meituan-batch-delivery-cost/"),
})

# ===================================== SLIDE 7 · HOW THE WORLD SOLVED IT
s = slide()
header(s, "07", "How The World Already Solved This", STEEL,
       P(R("Everyone shipped the cart. ", 26, INK, True, FONTD),
         R("Nobody shipped the coordination", 26, STEEL, True, FONTD)),
       P(R("Five food platforms converged on one weak mechanic — while four other industries already cracked the hard part. The blueprint isn't missing; it's just never been assembled for food.", 12, INK2)))

# --- left: competitive teardown table ---
text(s, 0.7, 2.06, 7.4, 0.24, [P(R("THE TEARDOWN · 6 PLATFORMS, ONE WEAK MECHANIC", 10, STEEL, True, FONT, 20))])
c0,w0, c1,w1, c2,w2, c3,w3 = 0.7,1.5, 2.2,1.6, 3.8,1.7, 5.5,2.6
ty, rowh = 2.4, 0.5
hd = [("PLATFORM",c0,w0),("REAL-TIME COORD.",c1,w1),("SPLIT-PAY",c2,w2),("THE GAP",c3,w3)]
for txt,cx,cw in hd:
    rect(s, cx, ty, cw, rowh, fill=STEEL, line=WHITE, line_w=1.0, radius=0)
    text(s, cx+0.08, ty, cw-0.14, rowh, [P(R(txt, 8.2, WHITE, True, FONT, 10))],
         anchor=MSO_ANCHOR.MIDDLE, line_spacing=0.9)
SY = {"✕":(CRIM,FAILBG),"~":(AMBER,GOLDBG),"✓":(PITCH,FIXBG),"—":(INK3,SOFT)}
trows = [
 ("Swiggy ’24","✕","none","✕","host pays all","Proved demand; skipped coordination"),
 ("Zomato ’24","✕","none","✕","“soon”, unshipped","At parity on the weak mechanic"),
 ("Uber Eats","~","limits","~","host eats fees","Controls, not coordination"),
 ("DoorDash","~","budgets","~","fees+tip on host","Best controls; fairness still broken"),
 ("Grab","~","QR / cutoff","✓","even split","Lone clean split — but no presence"),
 ("Meituan","✕","n/a","—","group-BUY","Discounts, not coordination"),
]
def ratecell(x,y,w,sym,word):
    col,tint = SY.get(sym,(INK3,SOFT))
    rect(s, x, y, w, rowh, fill=tint, line=WHITE, line_w=1.0, radius=0)
    text(s, x+0.1, y, w-0.16, rowh, [P(R(sym+" ", 12, col, True), R(word, 7.6, INK2))],
         anchor=MSO_ANCHOR.MIDDLE, line_spacing=0.88)
for ri,(plat,cs,cw_,ss,sw,verd) in enumerate(trows):
    ry = ty + (ri+1)*rowh
    rect(s, c0, ry, w0, rowh, fill=SOFT, line=WHITE, line_w=1.0, radius=0)
    text(s, c0+0.1, ry, w0-0.16, rowh, [P(R(plat, 9, INK, True))], anchor=MSO_ANCHOR.MIDDLE)
    ratecell(c1, ry, w1, cs, cw_)
    ratecell(c2, ry, w2, ss, sw)
    rect(s, c3, ry, w3, rowh, fill=WHITE, line=LINE, line_w=1.0, radius=0)
    text(s, c3+0.1, ry, w3-0.18, rowh, [P(R(verd, 8.3, INK2))], anchor=MSO_ANCHOR.MIDDLE, line_spacing=0.92)

# --- right: steal from outside food ---
text(s, 8.3, 2.06, 4.33, 0.24, [P(R("STEAL FROM OUTSIDE FOOD", 10, ZRED, True, FONT, 20))])
borrow = [
 ("Spotify Jam", VIOL, "32-person real-time queue; see who added what; the host runs the room.", "presence avatars"),
 ("Figma", STEEL, "Live cursors + per-person colour; no edit conflicts, ever.", "colour sub-carts"),
 ("Gaming lobbies", PITCH, "Everyone taps “ready”; host fires on a countdown; idle players dropped.", "ready-check + timer"),
 ("Splitwise / UPI", ZRED, "Everyone's share settled automatically — nobody chases ₹80s.", "clean split-pay"),
]
for i,(nm,col,desc,mech) in enumerate(borrow):
    by = 2.4 + i*0.96
    rect(s, 8.3, by, 4.33, 0.86, fill=WHITE, line=LINE, radius=0.05)
    avatar(s, 8.44, by+0.16, 0.34, "★", col)
    text(s, 8.9, by+0.1, 3.6, 0.26, [P(R(nm, 10.5, col, True, FONTD))])
    text(s, 8.9, by+0.34, 3.6, 0.3, [P(R(desc, 8.2, INK2))], line_spacing=0.95)
    text(s, 8.9, by+0.62, 3.6, 0.2, [P(R("→ borrow: ", 8.3, INK3, True), R(mech, 8.3, col, True))])

# cliffhanger
rect(s, 0.7, 6.4, 11.93, 0.6, fill=ZDEEP, line=None, radius=0.05)
text(s, 0.95, 6.4, 11.45, 0.6,
     [P(R("→  ", 14, RGBColor(0xFF,0xC9,0xC4), True),
        R("The patterns are already solved — just never assembled for food. ", 12, WHITE, True),
        R("That assembly is the pitch.", 12, RGBColor(0xFF,0xE0,0x9E), True))],
     anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0)
footer_links(s, srcs("grab","spotify_jam","doordash","meituan"), "07 / 14")

# ===================================== SLIDE 8 · THE PITCH · LIVE COLLABORATIVE CART
s = slide()
header(s, "08", "The Pitch · Live Collaborative Cart", PITCH,
       P(R("One live cart. ", 26, INK, True, FONTD), R("Every leak sealed.", 26, PITCH, True, FONTD)),
       P(R("A real-time, presence-rich, conflict-free room that does the coordinating for the group — each piece kills a specific cluster from the diagnosis.", 12, INK2)))

# concept band
rect(s, 0.7, 2.05, 11.93, 0.66, fill=FIXBG, line=FIXBD, radius=0.05)
text(s, 0.95, 2.05, 11.45, 0.66,
     [P(R("THE LIVE COLLABORATIVE CART   ", 12, PITCH, True, FONTD),
        R("host shares one live room → everyone adds in real time → the system drives consensus, timing, and fair payment.", 10, INK2))],
     anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0)

# 4 pillars
pillars = [
 ("Live presence","Ambient status — browsing · added · ready · paid. The room feels alive; “who's in?” just disappears.","KILLS COLLECTION", STEEL, STEELBG, "Spotify Jam"),
 ("Colour sub-carts","Each person's own cart, merged — never shared. Zero edit conflicts; every item attributed by colour.","KILLS COLLECTION", VIOL, RGBColor(0xF1,0xEC,0xF8), "Figma"),
 ("Ready-check","Each taps “I'm ready”; the host sees “3 of 5.” Consensus without nagging a soul.","KILLS TIMING", AMBER, GOLDBG, "Gaming lobbies"),
 ("Host timer","At the deadline: auto-fire confirmed items, drop stragglers, extend once. Nagging becomes a nudge.","KILLS TIMING", PITCH, FIXBG, "Gaming lobbies"),
]
for i,(nm,desc,kill,col,tint,src) in enumerate(pillars):
    px = 0.7 + i*(2.9+0.11)
    rect(s, px, 2.82, 2.9, 2.0, fill=tint, line=col, line_w=1.0, radius=0.05)
    rect(s, px, 2.82, 2.9, 0.46, fill=col, line=None, radius=0.05)
    text(s, px+0.12, 2.82, 2.66, 0.46, [P(R(nm, 12, WHITE, True, FONTD))], anchor=MSO_ANCHOR.MIDDLE)
    text(s, px+0.16, 3.4, 2.58, 0.94, [P(R(desc, 8.8, INK2))], line_spacing=1.02)
    chip(s, px+0.16, 4.32, 1.6, 0.28, kill, col)
    text(s, px+0.16, 4.64, 2.58, 0.2, [P(R("↳ from "+src, 8, INK3, True)) ])

# 2 bookends (discovery in, settlement out)
ends = [("Shareable deep-link","One tap shares the live room — discovery without digging through the cart.","DISCOVERY  ·  viral in", ZRED, ZBG),
        ("UPI-collect split","Each share auto-requested at checkout — the host stops being the group's bank.","SETTLEMENT  ·  fair out", ZDEEP, FAILBG)]
for i,(nm,desc,kill,col,tint) in enumerate(ends):
    bx = 0.7 + i*(5.86+0.21)
    rect(s, bx, 4.96, 5.86, 0.92, fill=tint, line=col, line_w=1.0, radius=0.05)
    text(s, bx+0.2, 5.04, 3.8, 0.3, [P(R(nm, 12, col, True, FONTD))])
    text(s, bx+0.2, 5.36, 4.0, 0.46, [P(R(desc, 9, INK2))], line_spacing=0.98)
    chip(s, bx+5.86-2.0, 5.08, 1.84, 0.3, "KILLS "+kill, col)

# every cluster sealed
rect(s, 0.7, 5.98, 11.93, 0.32, fill=FIXBG, line=FIXBD, radius=0.05)
text(s, 0.7, 5.98, 11.93, 0.32,
     [P(R("EVERY CLUSTER SEALED    ", 10, PITCH, True, FONT, 20),
        R("✓ Discovery    ✓ Collection    ✓ Settlement    ✓ Timing    ✓ Fulfilment", 10, INK, True))],
     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# cliffhanger
rect(s, 0.7, 6.42, 11.93, 0.58, fill=ZDEEP, line=None, radius=0.05)
text(s, 0.95, 6.42, 11.45, 0.58,
     [P(R("→  ", 14, RGBColor(0xFF,0xC9,0xC4), True),
        R("Five clusters, one cart, all sealed — but you can't ship it all on day one. ", 12, WHITE, True),
        R("So what's the v1?", 12, RGBColor(0xFF,0xE0,0x9E), True))],
     anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0)
footer_links(s, srcs("spotify_jam","grab","swiggy_blog"), "08 / 14")

VIOLBG = RGBColor(0xF1,0xEC,0xF8)
VIOLBD = RGBColor(0xD8,0xCA,0xEC)

# ===================================== SLIDE 9 · RICE & MVP SCOPE
s = slide()
header(s, "09", "Prioritisation · RICE & MVP Scope", VIOL,
       P(R("Six features, ", 26, INK, True, FONTD), R("one disciplined v1", 26, VIOL, True, FONTD)),
       P(R("RICE ranks the build; v1 ships only what moves the North Star and seeds the viral loop. Everything else is sequenced — not forgotten.", 12, INK2)))

text(s, 0.7, 2.02, 7.9, 0.24,
     [P(R("RICE SCORECARD   ", 10, VIOL, True, FONT, 20), R("Reach × Impact × Confidence ÷ Effort", 9, INK3, False, FONT, 0, True))])
f0,wf0, f1,wf1, f2,wf2, f3,wf3 = 0.7,2.25, 2.95,1.75, 4.7,1.05, 5.75,2.85
ty9, rh9 = 2.34, 0.55
for txt,cx,cw in [("FEATURE",f0,wf0),("RICE",f1,wf1),("STAGE",f2,wf2),("WHY IT RANKS HERE",f3,wf3)]:
    rect(s, cx, ty9, cw, rh9, fill=VIOL, line=WHITE, line_w=1.0, radius=0)
    text(s, cx+0.1, ty9, cw-0.16, rh9, [P(R(txt, 8.2, WHITE, True, FONT, 10))], anchor=MSO_ANCHOR.MIDDLE)
ricerows = [
 ("Collaborative cart",9.4,"v1",PITCH,"The only feature that moves the North Star"),
 ("WhatsApp deep-link",8.6,"v1",PITCH,"Low effort, huge reach — the viral surface"),
 ("UPI split-pay",7.1,"Stage 2",AMBER,"High value, high effort — the fairness fix"),
 ("Restaurant vote",6.2,"Stage 2",AMBER,"Cheap consensus, upstream of the cart"),
 ("Labelled sub-bags",4.5,"Stage 3",STEEL,"Delight — needs restaurant compliance"),
 ("Multi-restaurant",2.4,"Defer",INK3,"Compounds prep, routing & packaging"),
]
for ri,(nm,sc,stg,col,why) in enumerate(ricerows):
    ry = ty9 + (ri+1)*rh9
    rect(s, f0, ry, wf0, rh9, fill=SOFT, line=WHITE, line_w=1.0, radius=0)
    text(s, f0+0.1, ry, wf0-0.16, rh9, [P(R(nm, 9.3, INK, True))], anchor=MSO_ANCHOR.MIDDLE)
    rect(s, f1, ry, wf1, rh9, fill=WHITE, line=WHITE, line_w=1.0, radius=0)
    text(s, f1+0.08, ry, 0.5, rh9, [P(R(("%.1f"%sc), 11.5, col, True, FONTD))], anchor=MSO_ANCHOR.MIDDLE)
    rect(s, f1+0.62, ry+rh9/2-0.05, 1.02, 0.1, fill=RGBColor(0xEC,0xEF,0xF4), line=None, radius=0)
    rect(s, f1+0.62, ry+rh9/2-0.05, max(0.05,1.02*sc/10.0), 0.1, fill=col, line=None, radius=0)
    rect(s, f2, ry, wf2, rh9, fill=WHITE, line=WHITE, line_w=1.0, radius=0)
    chip(s, f2+0.06, ry+rh9/2-0.13, 0.92, 0.26, stg, col)
    rect(s, f3, ry, wf3, rh9, fill=WHITE, line=LINE, line_w=1.0, radius=0)
    text(s, f3+0.1, ry, wf3-0.18, rh9, [P(R(why, 8.2, INK2))], anchor=MSO_ANCHOR.MIDDLE, line_spacing=0.92)

# right: scope in / out
rect(s, 8.8, 2.34, 3.83, 2.04, fill=FIXBG, line=FIXBD, radius=0.05)
text(s, 8.98, 2.44, 3.5, 0.26, [P(R("SHIP NOW · v1", 10.5, PITCH, True, FONTD))])
text(s, 8.98, 2.74, 3.5, 1.6,
     [P(R("✓ ", 9.5, PITCH, True), R("Single restaurant + address", 9.2, INK)),
      P(R("✓ ", 9.5, PITCH, True), R("Host-pay default", 9.2, INK)),
      P(R("✓ ", 9.5, PITCH, True), R("Real-time cart: presence · sub-carts · ready-check · timer · nudge", 9.2, INK)),
      P(R("✓ ", 9.5, PITCH, True), R("No-install web cart for guests", 9.2, INK)),
      P(R("✓ ", 9.5, PITCH, True), R("Deep-link + QR + room code", 9.2, INK))],
     space_after=4, line_spacing=1.0)
rect(s, 8.8, 4.5, 3.83, 1.7, fill=SOFT, line=LINE, radius=0.05)
text(s, 8.98, 4.6, 3.5, 0.26, [P(R("SEQUENCED, NOT FORGOTTEN", 10.5, INK2, True, FONTD))])
text(s, 8.98, 4.9, 3.5, 1.3,
     [P(R("→ ", 9.5, AMBER, True), R("Split-pay · restaurant vote — Stage 2", 9.2, INK2)),
      P(R("→ ", 9.5, STEEL, True), R("Labelled sub-bags — Stage 3", 9.2, INK2)),
      P(R("→ ", 9.5, INK3, True), R("Multi-restaurant / multi-address — fast-follow", 9.2, INK2)),
      P(R("✕ ", 9.5, CRIM, True), R("Gold-gating — would kill the acquisition loop", 9.2, CRIM, True))],
     space_after=4, line_spacing=1.0)

rect(s, 0.7, 6.4, 11.93, 0.6, fill=ZDEEP, line=None, radius=0.05)
text(s, 0.95, 6.4, 11.45, 0.6,
     [P(R("→  ", 14, RGBColor(0xFF,0xC9,0xC4), True),
        R("v1 is the live cart and the link. ", 12, WHITE, True),
        R("Now — what does it actually look like?", 12, RGBColor(0xFF,0xE0,0x9E), True))],
     anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0)
footer_links(s, srcs("eternal_fy25","swiggy_li","spotify_jam"), "09 / 14")

# ===================================== SLIDE 10 · MOCKUP · HOST SCREEN
def phone(s, x, y, w, h):
    rect(s, x, y, w, h, fill=DARK, line=None, radius=0.14)
    rect(s, x+0.085, y+0.1, w-0.17, h-0.2, fill=WHITE, line=None, radius=0.11)
    o = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x+w/2-0.35), Inches(y+0.04), Inches(0.7), Inches(0.11))
    o.fill.solid(); o.fill.fore_color.rgb = RGBColor(0x2A,0x2F,0x3C); o.line.fill.background(); o.shadow.inherit=False
    _set_radius(o, 0.5)
    return (x+0.085, y+0.1, w-0.17, h-0.2)

def numbadge(s, x, y, d, n, col):
    o = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(d), Inches(d))
    o.fill.solid(); o.fill.fore_color.rgb = col; o.line.color.rgb = WHITE; o.line.width = Pt(1.25); o.shadow.inherit = False
    text(s, x, y, d, d, [P(R(str(n), 9.5, WHITE, True, FONTD))], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

def dot(s, x, y, d, col):
    o = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(d), Inches(d))
    o.fill.solid(); o.fill.fore_color.rgb = col; o.line.fill.background(); o.shadow.inherit=False

s = slide()
header(s, "10", "The Mockup · Host's Orchestration View", PITCH,
       P(R("What the host sees: ", 25, INK, True, FONTD), R("the whole room, one glance", 25, PITCH, True, FONTD)),
       P(R("This is the moat — not a shared cart, but a live control room. Every element on screen kills a coordination cost from the diagnosis.", 12, INK2)))

# ---- the phone (host orchestration screen) ----
sx, sy, sw, sh = phone(s, 0.7, 2.02, 3.05, 4.32)
ix = sx+0.14            # inner content x
iw = sw-0.28            # inner content w
# status bar
text(s, ix, sy+0.04, 1.0, 0.18, [P(R("9:41", 8, INK, True))], anchor=MSO_ANCHOR.MIDDLE)
text(s, sx+sw-1.0, sy+0.04, 0.86, 0.18, [P(R("▪ ▪ ▪  ▰", 7.5, INK2))], align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
# title row
text(s, ix, sy+0.26, 1.8, 0.28, [P(R("Group Order", 12.5, INK, True, FONTD))], anchor=MSO_ANCHOR.MIDDLE)
o = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(sx+sw-1.0), Inches(sy+0.28), Inches(0.86), Inches(0.26))
o.fill.solid(); o.fill.fore_color.rgb=ZRED; o.line.fill.background(); o.shadow.inherit=False; _set_radius(o,0.45)
dot(s, sx+sw-0.9, sy+0.36, 0.1, WHITE)
text(s, sx+sw-1.0, sy+0.28, 0.86, 0.26, [P(R("LIVE", 8.5, WHITE, True, FONT, 30))], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
text(s, ix, sy+0.56, 2.6, 0.2, [P(R("Meghana Foods · 2.1 km", 8, INK2))])
bar(s, ix, sy+0.82, iw, 0.01, LINE)
# presence avatars
av = [("R",ZRED),("P",STEEL),("A",VIOL),("S",AMBER),("+1",INK3)]
for i,(ini,col) in enumerate(av):
    avatar(s, ix+i*0.5, sy+0.92, 0.42, ini, col)
text(s, ix+2.55, sy+0.98, 0.9, 0.3, [P(R("5 joined", 8, INK2, True))], anchor=MSO_ANCHOR.MIDDLE)
# sub-cart cards
cards = [("R",ZRED,"Rahul","Chicken Biryani ×1","₹320","Ready",PITCH),
         ("P",STEEL,"Priya","Paneer Masala · Naan","₹350","Ready",PITCH),
         ("A",VIOL,"Aman","Chicken 65","₹260","Browsing",AMBER)]
for i,(ini,col,nm,item,price,stt,scol) in enumerate(cards):
    cy = sy+1.44 + i*0.54
    rect(s, ix, cy, iw, 0.48, fill=WHITE, line=LINE, radius=0.07)
    dot(s, ix+0.12, cy+0.09, 0.18, col)
    text(s, ix+0.4, cy+0.05, 1.5, 0.22, [P(R(nm, 8.6, INK, True))])
    text(s, ix+0.4, cy+0.25, 1.7, 0.2, [P(R(item, 7.2, INK2))])
    text(s, ix+iw-1.0, cy+0.05, 0.92, 0.22, [P(R(price, 9, INK, True, FONTD))], align=PP_ALIGN.RIGHT)
    o = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(ix+iw-1.0), Inches(cy+0.27), Inches(0.92), Inches(0.18))
    o.fill.solid(); o.fill.fore_color.rgb = (FIXBG if stt=="Ready" else GOLDBG); o.line.fill.background(); o.shadow.inherit=False; _set_radius(o,0.5)
    text(s, ix+iw-1.0, cy+0.255, 0.92, 0.2, [P(R(stt, 6.8, scol, True)) ], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
# totals card
rect(s, ix, sy+3.06, iw, 0.66, fill=ZBG, line=ZBD, radius=0.07)
text(s, ix+0.12, sy+3.12, 2.6, 0.18, [P(R("Group total · MOV ₹199 cleared ✓", 7, PITCH, True))])
text(s, ix+0.12, sy+3.3, 1.6, 0.34, [P(R("₹1,240", 16, INK, True, FONTD))])
chip(s, ix+iw-1.25, sy+3.16, 1.13, 0.26, "3 of 5 ready", PITCH)
text(s, ix+iw-1.25, sy+3.46, 1.13, 0.2, [P(R("⏱ 04:12 left", 8, ZDEEP, True))], align=PP_ALIGN.RIGHT)
# CTA (waiting state)
o = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(ix), Inches(sy+3.8), Inches(iw), Inches(0.44))
o.fill.solid(); o.fill.fore_color.rgb=ZRED; o.line.fill.background(); o.shadow.inherit=False; _set_radius(o,0.18)
text(s, ix, sy+3.8, iw, 0.44, [P(R("Waiting on 2  ·  Nudge or place confirmed  →", 9, WHITE, True, FONTD))],
     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# ---- right: annotation rows ----
text(s, 4.15, 2.12, 8.45, 0.26, [P(R("EVERY PIXEL KILLS A COORDINATION COST", 11, ZDEEP, True, FONT, 20))])
bar(s, 4.15, 2.4, 8.45, 0.02, ZBD)
calls = [
 (1,"Live presence avatars",STEEL,"Browsing · ready · paid at a glance — the room finally feels alive.","COLLECTION"),
 (2,"Colour sub-carts",VIOL,"Every item owned and attributed — no mix-ups, no edit conflicts, ever.","COLLECTION"),
 (3,"Ambient status chips",AMBER,"Ready / browsing / away — calm presence, never the anxiety of a shared list.","TIMING"),
 (4,"Ready-check + live total",PITCH,"“3 of 5 ready” and MOV auto-cleared — consensus without nagging a soul.","TIMING"),
 (5,"Host timer",ZRED,"⏱ Auto-fires at 0:00; drop stragglers, extend once. Nagging becomes a nudge.","TIMING"),
 (6,"Smart place button",ZDEEP,"“Waiting on 2 → place 3 confirmed.” The button's state drives the action.","SETTLEMENT"),
]
ry0, rhh = 2.52, 0.63
for i,(n,t,col,desc,clu) in enumerate(calls):
    cy = ry0 + i*rhh
    rect(s, 4.15, cy, 8.45, 0.57, fill=WHITE, line=LINE, radius=0.05)
    numbadge(s, 4.29, cy+0.14, 0.3, n, col)
    text(s, 4.74, cy+0.07, 4.2, 0.24, [P(R(t, 10.5, col, True, FONTD))])
    text(s, 4.74, cy+0.31, 6.4, 0.22, [P(R(desc, 8.5, INK2))])
    chip(s, 4.15+8.45-1.42, cy+0.16, 1.28, 0.26, clu, col)

rect(s, 0.7, 6.46, 11.93, 0.52, fill=ZDEEP, line=None, radius=0.05)
text(s, 0.95, 6.46, 11.45, 0.52,
     [P(R("→  ", 13, RGBColor(0xFF,0xC9,0xC4), True),
        R("That's the host's orchestration view. Next: the guest's one-tap web cart — ", 11.5, WHITE, True),
        R("and the bill that splits itself.", 11.5, RGBColor(0xFF,0xE0,0x9E), True))],
     anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0)
footer_links(s, srcs("spotify_jam","grab","doordash"), "10 / 14")

def mockspace(s, x, y, w, h, label, sub):
    shp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    _set_radius(shp, 0.05)
    shp.fill.solid(); shp.fill.fore_color.rgb = RGBColor(0xF3,0xF5,0xF9)
    shp.line.color.rgb = INK3; shp.line.width = Pt(1.25)
    try:
        from pptx.enum.line import MSO_LINE_DASH_STYLE
        shp.line.dash_style = MSO_LINE_DASH_STYLE.DASH
    except Exception: pass
    shp.shadow.inherit = False
    bar(s, x+0.22, y+0.34, w-0.44, 0.42, RGBColor(0xDF,0xE4,0xED))            # header hint
    rect(s, x+0.22, y+h-0.72, w-0.44, 0.46, fill=RGBColor(0xD6,0xDC,0xE6), line=None, radius=0.12)  # CTA hint
    text(s, x, y+h/2-0.42, w, 0.9,
         [P(R("▭", 22, INK3, True)), P(R(label, 11, INK2, True, FONTD)),
          P(R(sub, 8.3, INK3)), P(R("mockup to be placed", 7.5, INK3, False, FONT, 0, True))],
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, space_after=2, line_spacing=1.0)

TINT = {ZRED:ZBG, AMBER:GOLDBG, VIOL:VIOLBG, PITCH:FIXBG, STEEL:STEELBG, ZDEEP:FAILBG}

# ===================================== SLIDE 11 · MOCKUP · GUEST WEB CART
s = slide()
header(s, "11", "The Mockup · Guest's One-Tap Cart", PITCH,
       P(R("No install. One tap. ", 25, INK, True, FONTD), R("they're in.", 25, PITCH, True, FONTD)),
       P(R("The invitee never downloads a thing. A link opens a live web cart — they add, mark ready, and pay only their share. Friction out, fairness in.", 12, INK2)))

mockspace(s, 0.7, 2.05, 3.05, 4.3, "INVITEE WEB CART", "guest screen")

text(s, 4.15, 2.12, 8.45, 0.26, [P(R("WHAT THE GUEST EXPERIENCES", 11, PITCH, True, FONT, 20))])
bar(s, 4.15, 2.4, 8.45, 0.02, FIXBD)
grows = [
 (1,"“You're in Rahul's order · WEB”",STEEL,"No app, no login — the link opens a live cart right in the browser.","DISCOVERY"),
 (2,"Browse & tap +",VIOL,"Add from the full menu, exactly like a normal solo order.","COLLECTION"),
 (3,"“Your sub-cart”",AMBER,"Only your items and your subtotal — sticky at the bottom.","COLLECTION"),
 (4,"“I'm ready ✓”",PITCH,"One tap feeds the host's ready-check. No more “you done yet?”","TIMING"),
 (5,"“Pay my ₹392 via UPI →”",ZRED,"Your share only — the host finally stops being the group's bank.","SETTLEMENT"),
]
gy0 = 2.5
for i,(n,t,col,desc,clu) in enumerate(grows):
    cy = gy0 + i*0.6
    rect(s, 4.15, cy, 8.45, 0.55, fill=WHITE, line=LINE, radius=0.05)
    numbadge(s, 4.29, cy+0.13, 0.3, n, col)
    text(s, 4.74, cy+0.06, 4.6, 0.24, [P(R(t, 10.5, col, True, FONTD))])
    text(s, 4.74, cy+0.3, 6.4, 0.22, [P(R(desc, 8.5, INK2))])
    chip(s, 4.15+8.45-1.42, cy+0.15, 1.28, 0.26, clu, col)

# strategic payoff
payoff = [("VIRAL ACQUISITION", ZRED, ZBG, ZBD, "Every guest who pays is a non-user touched at zero CAC — the K-factor loop."),
          ("SETTLEMENT SOLVED", PITCH, FIXBG, FIXBD, "UPI-collect ends the ₹80 chase. The single most-requested fix, shipped.")]
for i,(lab,col,fb,bd,txt) in enumerate(payoff):
    px = 4.15 + i*4.28
    rect(s, px, 5.56, 4.1, 0.62, fill=fb, line=bd, radius=0.05)
    text(s, px+0.16, 5.6, 3.8, 0.54, [P(R(lab+"  ", 9, col, True, FONTD), R(txt, 8.3, INK2))],
         anchor=MSO_ANCHOR.MIDDLE, line_spacing=0.96)

rect(s, 0.7, 6.42, 11.93, 0.56, fill=ZDEEP, line=None, radius=0.05)
text(s, 0.95, 6.42, 11.45, 0.56,
     [P(R("→  ", 13, RGBColor(0xFF,0xC9,0xC4), True),
        R("Beautiful on the happy path. But real groups are messy — ", 11.5, WHITE, True),
        R("what happens when it breaks?", 11.5, RGBColor(0xFF,0xE0,0x9E), True))],
     anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0)
footer_links(s, srcs("grab","ubereats","doordash"), "11 / 14")

# ===================================== SLIDE 12 · EDGE CASES
s = slide()
header(s, "12", "Built To Not Break · Edge Cases", STEEL,
       P(R("Real groups are messy. ", 25, INK, True, FONTD), R("The system isn't.", 25, STEEL, True, FONTD)),
       P(R("Happy paths are easy; the moat is in the edge cases. Every way a group order can fall apart — each with a defined, no-host-burden system response.", 12, INK2)))

mockspace(s, 0.7, 2.05, 3.05, 4.3, "THE ORDER IN MOTION", "browsing → ready → fired")

text(s, 4.15, 2.12, 8.45, 0.26, [P(R("EVERY WAY IT BREAKS — AND EXACTLY WHAT HAPPENS", 11, STEEL, True, FONT, 20))])
bar(s, 4.15, 2.4, 8.45, 0.02, STEELBD)
edges = [
 ("Item sells out", ZRED, "Real-time flag tagged to that member; blocks ready-check until they re-pick."),
 ("Member goes idle", AMBER, "Presence dims to “away”; host can poke once, or drop them."),
 ("Host abandons", VIOL, "Transfer-host to the longest-active member — the session never dies."),
 ("Timer expires", PITCH, "Auto-fire confirmed items; drop unconfirmed (notified); host can extend once."),
 ("Late joiner", STEEL, "Cart locks at checkout; they see “ordering started” and can start a fresh one."),
 ("Member never pays", ZDEEP, "Host is prompted to nudge or drop the items — never silently charged."),
]
ey0, erh = 2.52, 0.6
for i,(edge,col,resp) in enumerate(edges):
    cy = ey0 + i*erh
    rect(s, 4.15, cy, 2.75, 0.54, fill=TINT[col], line=WHITE, line_w=1.0, radius=0.05)
    dot(s, 4.29, cy+0.19, 0.16, col)
    text(s, 4.54, cy, 2.3, 0.54, [P(R(edge, 9.3, col, True, FONTD))], anchor=MSO_ANCHOR.MIDDLE)
    rect(s, 6.95, cy, 5.65, 0.54, fill=WHITE, line=LINE, line_w=1.0, radius=0.05)
    text(s, 7.1, cy, 5.4, 0.54, [P(R(resp, 8.8, INK2))], anchor=MSO_ANCHOR.MIDDLE, line_spacing=0.95)

rect(s, 0.7, 6.42, 11.93, 0.56, fill=ZDEEP, line=None, radius=0.05)
text(s, 0.95, 6.42, 11.45, 0.56,
     [P(R("→  ", 13, RGBColor(0xFF,0xC9,0xC4), True),
        R("Robust, fair and delightful. ", 11.5, WHITE, True),
        R("Now the only question left — is it worth building?", 11.5, RGBColor(0xFF,0xE0,0x9E), True))],
     anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0)
footer_links(s, srcs("ubereats","doordash","spotify_jam"), "12 / 14")

# ===================================== SLIDE 13 · THE BUSINESS CASE
s = slide()
header(s, "13", "The Business Case · Does It Pay?", ZRED,
       P(R("Conservative upside: ", 25, INK, True, FONTD), R("₹560 crore GMV. Aggressive: ₹2,560 crore.", 25, ZRED, True, FONTD)),
       P(R("Three assumptions unlock outsized returns — higher AOV on group orders, modest lift in group-session frequency, and a zero-CAC viral loop from guest invites.", 12, INK2)))

# ---- baseline strip ----
rect(s, 0.7, 2.06, 11.93, 0.46, fill=SOFT, line=LINE, radius=0.05)
baselines = [("853M","orders / yr  (FY25)"),("₹453","solo AOV"),("₹38,640 cr","GMV (FY25)"),("20.6M","MTU")]
for i,(num,lbl) in enumerate(baselines):
    bx = 0.95 + i*3.0
    text(s, bx, 2.06, 2.8, 0.46,
         [P(R(num+"  ", 14, ZDEEP, True, FONTD), R(lbl, 9, INK2))],
         anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0)
text(s, 0.7, 2.06, 11.93, 0.46, [P(R("FY25 BASELINE  →", 8.5, ZRED, True, FONT, 10))],
     anchor=MSO_ANCHOR.MIDDLE)

# ---- 3 scenario cards ----
scenarios = [
 ("Conservative","₹560 cr","₹112 cr net","6% group-order adoption · +35% AOV lift",STEEL,STEELBG,STEELBD),
 ("Base Case","₹1,177 cr","₹235 cr net","12% adoption · +38% AOV · K-factor 0.15",ZRED,ZBG,ZBD),
 ("Aggressive","₹2,560 cr","₹512 cr net","18% adoption · +40% AOV · K-factor 0.30+",PITCH,FIXBG,FIXBD),
]
for i,(lab,gmv,net,assum,col,fbg,fbd) in enumerate(scenarios):
    cx = 0.7 + i*3.92
    rect(s, cx, 2.66, 3.75, 2.3, fill=fbg, line=fbd, radius=0.06)
    rect(s, cx, 2.66, 3.75, 0.52, fill=col, line=None, radius=0.06)
    text(s, cx+0.16, 2.66, 3.43, 0.52, [P(R(lab.upper(), 11, WHITE, True, FONTD))], anchor=MSO_ANCHOR.MIDDLE)
    text(s, cx+0.16, 3.28, 3.43, 0.56,
         [P(R(gmv, 26, col, True, FONTD))], anchor=MSO_ANCHOR.MIDDLE)
    chip(s, cx+0.16, 3.9, 1.8, 0.28, net, col)
    text(s, cx+0.16, 4.28, 3.43, 0.6, [P(R(assum, 8.2, INK2))], line_spacing=1.0)

# ---- unit economics strip ----
rect(s, 0.7, 5.1, 11.93, 0.48, fill=ZBG, line=ZBD, radius=0.05)
econs = [("~₹3–5 cr","build cost"),("47×","ROI · base case"),("~16 days","payback"),("40×","safety margin on breakeven")]
for i,(num,lbl) in enumerate(econs):
    ex = 0.95 + i*2.98
    text(s, ex, 5.1, 2.8, 0.48,
         [P(R(num+"  ", 13, ZDEEP, True, FONTD), R(lbl, 8.5, INK2))],
         anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0)

# ---- guardrails ----
rect(s, 0.7, 5.7, 11.93, 0.56, fill=FAILBG, line=FAILBD, radius=0.05)
text(s, 0.95, 5.7, 11.45, 0.56,
     [P(R("GUARDRAILS  ·  ", 9.5, CRIM, True, FONT, 15),
        R("cannibalization rate  ·  dispute/refund rate  ·  host-initiated cancellations  ·  support contacts per group order", 9, INK2))],
     anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0)

rect(s, 0.7, 6.42, 11.93, 0.56, fill=ZDEEP, line=None, radius=0.05)
text(s, 0.95, 6.42, 11.45, 0.56,
     [P(R("→  ", 13, RGBColor(0xFF,0xC9,0xC4), True),
        R("The unit economics make this a near-zero-risk build. ", 11.5, WHITE, True),
        R("The real prize is structural — it builds a moat.", 11.5, RGBColor(0xFF,0xE0,0x9E), True))],
     anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0)
footer_links(s, srcs("eternal_fy25","eternal_pdf","pendo"), "13 / 14")

# ===================================== SLIDE 14 · THE PRIZE · CLOSE
s = slide()
header(s, "14", "The Prize · Why This Builds A Moat", ZRED,
       P(R("Three compounding levers. ", 25, INK, True, FONTD), R("One social graph nobody else has.", 25, ZRED, True, FONTD)),
       P(R("Group ordering isn't a feature — it's a wedge into the social layer of food. Each lever compounds the others; together they lock in a network effect Swiggy can't copy without users.", 12, INK2)))

# 3 levers
levers = [
 ("01","Higher AOV","Group orders drive basket sizes 35–40% above solo — more items, add-ons and drinks from the host, every single session.",
  ZRED, ZBG, ZBD, "₹453 → ₹610+ average"),
 ("02","Higher Frequency","Solving the coordination tax removes the #1 drop-off reason for group sessions — groups order more often because it's no longer painful.",
  VIOL, VIOLBG, VIOLBD, "+2–4 sessions / month / group"),
 ("03","Viral K-factor","Every host invite is a zero-CAC acquisition touch. At K ≥ 0.3, each group order organically acquires 0.3 new users into the funnel.",
  PITCH, FIXBG, FIXBD, "K > 0.3 → viral coefficient"),
]
for i,(num,nm,desc,col,fbg,fbd,kpi) in enumerate(levers):
    lx = 0.7 + i*3.98
    rect(s, lx, 2.12, 3.8, 2.56, fill=fbg, line=fbd, radius=0.07)
    rect(s, lx, 2.12, 3.8, 0.56, fill=col, line=None, radius=0.07)
    text(s, lx+0.2, 2.12, 0.46, 0.56, [P(R(num, 9, WHITE, True, FONTD))], anchor=MSO_ANCHOR.MIDDLE)
    text(s, lx+0.7, 2.12, 2.9, 0.56, [P(R(nm, 13, WHITE, True, FONTD))], anchor=MSO_ANCHOR.MIDDLE)
    text(s, lx+0.2, 2.76, 3.4, 1.14, [P(R(desc, 9.2, INK2))], line_spacing=1.04)
    chip(s, lx+0.2, 4.12, 3.4, 0.3, kpi, col)

# moat statement
rect(s, 0.7, 4.86, 11.93, 0.82, fill=ZBG, line=ZBD, radius=0.07)
text(s, 0.95, 4.86, 11.45, 0.82,
     [P(R("THE MOAT  ·  ", 10.5, ZRED, True, FONTD),
        R("The moment a group places its first order on Zomato, their food-social graph lives here — ", 10, INK2),
        R("not on Swiggy, not on WhatsApp.", 10, ZDEEP, True),
        R("  Switching means re-inviting the whole group, rebuilding trust, and losing split history. That's structural lock-in, and it compounds every session.", 10, INK2))],
     anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.06)

# final CTA
rect(s, 0.7, 5.82, 11.93, 1.16, fill=ZDEEP, line=None, radius=0.07)
text(s, 0.95, 5.82, 11.45, 1.16,
     [P(R("The live cart isn't a feature drop. It's Zomato becoming the social infrastructure of group food.", 14, WHITE, True, FONTD)),
      P(R("The coordination cost is ours to eliminate — and the social graph is ours to own.", 11.5, RGBColor(0xFF,0xC9,0xC4)))],
     anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.12, space_after=6)
footer_links(s, srcs("eternal_fy25","pendo","swiggy_li"), "14 / 14")

out = safe_save(prs, "Zomato-Group-Ordering-PM-Deck.pptx")
print("saved", out, " ·  slides:", len(prs.slides._sldIdLst))
